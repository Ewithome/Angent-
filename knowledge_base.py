"""企业级内部知识库：多格式解析、语义分块、向量检索与混合检索。"""
from __future__ import annotations

import os
import re
from pathlib import Path

from langchain.tools import tool
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

KNOWLEDGE_DIR = Path(os.getenv("KNOWLEDGE_DIR", "knowledge"))
_SUPPORTED_SUFFIXES = {".md", ".txt", ".pdf", ".docx", ".pptx"}
_CHUNK_SIZE = 800
_CHUNK_OVERLAP = 120
_TOP_K = 3
_VECTOR_WEIGHT = 0.4  # 向量相似度在混合检索中的权重

_index_cache: dict = {"key": None, "chunks": [], "vectorizer": None, "matrix": None}


def _clean_text(text: str) -> str:
    """文档清洗：折叠空白、去掉孤立页码和无意义符号，保留中文标点。"""
    text = re.sub(r"[ \t\u3000]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"^\s*(第\s*\d+\s*页|PAGE\s*\d+)\s*$", "", text, flags=re.MULTILINE | re.IGNORECASE)
    return text.strip()


def _chunk_text(text: str, size: int = _CHUNK_SIZE, overlap: int = _CHUNK_OVERLAP) -> list[str]:
    """按固定窗口切分文本，重叠部分避免切断语义。"""
    text = _clean_text(text)
    if len(text) <= size:
        return [text] if text else []

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + size, len(text))
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = end - overlap
    return chunks


def _extract_text(path: Path) -> list[tuple[str, str]]:
    """按文件类型提取文本，返回 (来源描述, 文本) 列表。"""
    suffix = path.suffix.lower()
    if suffix in {".md", ".txt"}:
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = path.read_text(encoding="gbk", errors="ignore")
        return [(path.name, text)]

    if suffix == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = []
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append((f"{path.name} 第{index}页", text))
        return pages

    if suffix == ".docx":
        from docx import Document

        doc = Document(str(path))
        parts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        text = "\n".join(parts)
        return [(path.name, text)]

    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(str(path))
        slides = []
        for index, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    parts.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            parts.append(" | ".join(cells))
            if parts:
                slides.append(f"第{index}页：\n" + "\n".join(parts))
        return [(path.name, "\n\n".join(slides))]

    return []


def _index_key() -> tuple:
    """用文件路径、修改时间和大小作为索引缓存 key。"""
    if not KNOWLEDGE_DIR.exists():
        return ()
    return tuple(
        (
            str(path.relative_to(KNOWLEDGE_DIR)),
            path.stat().st_mtime_ns,
            path.stat().st_size,
        )
        for path in sorted(KNOWLEDGE_DIR.rglob("*"))
        if (
            path.is_file()
            and path.suffix.lower() in _SUPPORTED_SUFFIXES
            and path.name.lower() != "readme.md"
        )
    )


def _build_index() -> dict:
    """惰性构建文档分块与向量索引，文件变化后自动重建。"""
    key = _index_key()
    if key == _index_cache["key"]:
        return _index_cache

    chunks: list[dict] = []
    for path in sorted(KNOWLEDGE_DIR.rglob("*")):
        if (
            not path.is_file()
            or path.suffix.lower() not in _SUPPORTED_SUFFIXES
            or path.name.lower() == "readme.md"
        ):
            continue
        for source, text in _extract_text(path):
            for part in _chunk_text(text):
                if part:
                    chunks.append({"source": source, "text": part})

    vectorizer = None
    matrix = None
    if chunks:
        # 轻量向量检索：TF-IDF 字符级 n-gram，适合中文；后续可替换为 BGE 等语义嵌入模型
        vectorizer = TfidfVectorizer(
            analyzer="char_wb",
            ngram_range=(1, 2),
            max_features=50000,
            sublinear_tf=True,
        )
        matrix = vectorizer.fit_transform([chunk["text"] for chunk in chunks])

    _index_cache.update(
        {
            "key": key,
            "chunks": chunks,
            "vectorizer": vectorizer,
            "matrix": matrix,
        }
    )
    return _index_cache


def _tokenize(text: str) -> list[str]:
    """英文按单词、中文按单字切分，适合规范条文检索。"""
    return re.findall(r"[a-zA-Z0-9]+|[\u4e00-\u9fff]", text.lower())


def get_knowledge_file_count() -> int:
    """返回知识库中的文档数量，用于页面状态展示。"""
    if not KNOWLEDGE_DIR.exists():
        return 0
    return sum(
        1
        for path in KNOWLEDGE_DIR.rglob("*")
        if (
            path.is_file()
            and path.suffix.lower() in _SUPPORTED_SUFFIXES
            and path.name.lower() != "readme.md"
        )
    )


def build_knowledge_index() -> tuple[int, int]:
    """构建索引并返回 (文档数, 分块数)，供后台脚本与评测使用。"""
    cache = _build_index()
    return get_knowledge_file_count(), len(cache["chunks"])


def list_chunks() -> list[dict]:
    """返回当前索引的全部语义分块，供伪标签生成等离线任务使用。"""
    return _build_index()["chunks"]


def retrieve(query: str, top_k: int = 5) -> list[dict]:
    """混合检索：BM25 关键词得分 + TF-IDF 向量余弦相似度。"""
    cache = _build_index()
    chunks = cache["chunks"]
    if not chunks:
        return []

    query_terms = _tokenize(query)
    bm25_scores = [
        sum(chunk["text"].lower().count(term) for term in query_terms)
        for chunk in chunks
    ]

    vector_scores = [0.0] * len(chunks)
    if cache["vectorizer"] is not None and cache["matrix"] is not None:
        query_vector = cache["vectorizer"].transform([query])
        vector_scores = cosine_similarity(cache["matrix"], query_vector).ravel()

    max_bm25 = max(bm25_scores) or 1
    scored = []
    for index, chunk in enumerate(chunks):
        bm25_norm = bm25_scores[index] / max_bm25
        combined = bm25_norm * (1 - _VECTOR_WEIGHT) + vector_scores[index] * _VECTOR_WEIGHT
        scored.append(
            (
                combined,
                bm25_scores[index],
                vector_scores[index],
                chunk,
            )
        )

    scored.sort(key=lambda item: item[0], reverse=True)
    return [
        {
            "source": item[3]["source"],
            "text": item[3]["text"],
            "score": round(item[0], 4),
            "bm25_score": item[1],
            "vector_score": round(item[2], 4),
        }
        for item in scored[:top_k]
    ]


@tool
def search_building_code(query: str) -> str:
    """从企业内部知识库检索相关文档内容。参数 query 是问题或关键词，例如“楼梯踏步高度”或“产品退换货流程”。"""
    results = retrieve(query, top_k=_TOP_K)
    if not results or results[0]["score"] <= 0:
        return (
            "知识库为空或未找到相关内容。请把 PDF、Word、PPT、Markdown 或 TXT 文件放入 knowledge/ 目录后重试。"
            "在没有规范依据时，不要编造条文编号或强制条款。"
        )

    lines = [f"找到以下相关内容（共 {len(results)} 条）："]
    for item in results:
        snippet = item["text"][:600].replace("\n", " ")
        lines.append(
            f"- 来源：{item['source']}（综合得分 {item['score']}）\n  {snippet}"
        )
    return "\n\n".join(lines)
